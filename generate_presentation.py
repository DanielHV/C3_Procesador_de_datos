"""
Genera una sola diapositiva con el diagrama del pipeline C3_Procesador_de_datos.
Requiere: pip install python-pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

AZUL_OSC  = RGBColor(0x0D, 0x2B, 0x4E)
AZUL_MED  = RGBColor(0x1A, 0x52, 0x76)
AZUL_CLA  = RGBColor(0x2E, 0x86, 0xC1)
NARANJA   = RGBColor(0xE6, 0x7E, 0x22)
GRIS      = RGBColor(0xF2, 0xF4, 0xF4)
BLANCO    = RGBColor(0xFF, 0xFF, 0xFF)
NEGRO     = RGBColor(0x1C, 0x1C, 0x1C)

W = Inches(13.33)
H = Inches(7.5)


def rect(slide, l, t, w, h, fill=None, line=None):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid() if fill else s.fill.background()
    if fill:
        s.fill.fore_color.rgb = fill
    s.line.fill.background() if not line else None
    if line:
        s.line.color.rgb = line
    return s


def txt(slide, text, l, t, w, h, size=14, bold=False,
        color=NEGRO, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text = text
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.color.rgb = color
    return tb


def arrow_h(slide, x, y, length, color=AZUL_MED):
    """Flecha horizontal simple."""
    rect(slide, x, y + Inches(0.07), length - Inches(0.1), Inches(0.06), fill=color)
    # punta
    tip_w = Inches(0.18)
    tip_h = Inches(0.2)
    tip = slide.shapes.add_shape(
        5,  # triángulo isósceles
        x + length - tip_w, y - tip_h / 2 + Inches(0.1),
        tip_w, tip_h
    )
    tip.fill.solid()
    tip.fill.fore_color.rgb = color
    tip.line.fill.background()
    # rotar 90° para que apunte a la derecha
    tip.rotation = 90


def arrow_v(slide, x, y, length, color=AZUL_MED):
    """Flecha vertical hacia abajo (solo línea)."""
    rect(slide, x + Inches(0.07), y, Inches(0.06), length, fill=color)


def box(slide, label, sublabel, l, t, w, h, fill, text_color=BLANCO):
    rect(slide, l, t, w, h, fill=fill)
    # label centrado verticalmente
    label_h = Inches(0.4)
    label_t = t + (h - label_h) / 2 - Inches(0.15)
    txt(slide, label, l + Inches(0.1), label_t, w - Inches(0.2),
        label_h, size=13, bold=True, color=text_color, align=PP_ALIGN.CENTER)
    if sublabel:
        txt(slide, sublabel, l + Inches(0.1), label_t + Inches(0.38),
            w - Inches(0.2), Inches(0.35), size=10,
            color=text_color, align=PP_ALIGN.CENTER)


prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

slide = prs.slides.add_slide(prs.slide_layouts[6])

# ── Fondo ──────────────────────────────────────────────────────────────────
rect(slide, 0, 0, W, H, fill=GRIS)

# ── Título ─────────────────────────────────────────────────────────────────
rect(slide, 0, 0, W, Inches(0.85), fill=AZUL_OSC)
rect(slide, 0, 0, Inches(0.2), H, fill=NARANJA)
txt(slide, "C3 Procesador de Datos  —  Flujo del Pipeline",
    Inches(0.4), Inches(0.18), Inches(12.5), Inches(0.6),
    size=22, bold=True, color=BLANCO)

# ══════════════════════════════════════════════════════════════════════════════
# NODOS principales (izquierda → derecha)
# Coordenadas: y centrado en la parte alta del diagrama
# ══════════════════════════════════════════════════════════════════════════════
BW = Inches(2.2)   # ancho de cada caja
BH = Inches(1.1)   # alto de cada caja
BY = Inches(2.5)   # y de las cajas principales
GAP = Inches(0.7)  # espacio entre cajas

X0 = Inches(0.45)  # inicio de la primera caja

nodes = [
    ("Datos crudos",       "CSV + Diccionario",  AZUL_CLA),
    ("Preprocesamiento",   "Agrupación & alias", AZUL_MED),
    ("Procesamiento",      "Cuantiles & presencia", AZUL_OSC),
    ("Carga a BD",         "PostgreSQL",         NARANJA),
]

xs = [X0 + i * (BW + GAP) for i in range(len(nodes))]

for i, (label, sub, color) in enumerate(nodes):
    box(slide, label, sub, xs[i], BY, BW, BH, fill=color)

# Flechas horizontales entre cajas
arr_y = BY + BH / 2
for i in range(len(nodes) - 1):
    x_start = xs[i] + BW
    x_end   = xs[i + 1]
    length  = x_end - x_start
    arrow_h(slide, x_start, arr_y - Inches(0.03), length)

# ══════════════════════════════════════════════════════════════════════════════
# DETALLE de cada nodo (cajas secundarias debajo)
# ══════════════════════════════════════════════════════════════════════════════
DH  = Inches(0.38)
DY0 = BY + BH + Inches(0.6)   # y del primer detalle
D_COLS = [
    [   # Datos crudos
        ("Archivo de datos (.csv)", AZUL_CLA),
        ("Diccionario de variables", AZUL_CLA),
        ("Config JSON (rutas, columnas)", AZUL_CLA),
    ],
    [   # Preprocesamiento
        ("Limpieza y cast de tipos", AZUL_MED),
        ("Agrupación categórica → conteos", AZUL_MED),
        ("Agrupación numérica → media/mediana", AZUL_MED),
    ],
    [   # Procesamiento
        ("Normalización por variable base", AZUL_OSC),
        ("Clasificación en q cuantiles", AZUL_OSC),
        ("Detección de presencia (≥ 1)", AZUL_OSC),
    ],
    [   # BD
        ("Tabla mallas (_mun, _estado…)", NARANJA),
        ("Tabla ensamble (_personas/_empresas)", NARANJA),
        ("Diccionario y FK (dict_, values_)", NARANJA),
    ],
]

ALPHA = {  # versión más clara para fondo de sub-cajas
    AZUL_CLA: RGBColor(0xD6, 0xE9, 0xF8),
    AZUL_MED: RGBColor(0xC8, 0xDB, 0xEA),
    AZUL_OSC: RGBColor(0xBD, 0xCB, 0xD9),
    NARANJA:  RGBColor(0xFA, 0xDF, 0xC6),
}
TEXT_C = {
    AZUL_CLA: AZUL_OSC,
    AZUL_MED: AZUL_OSC,
    AZUL_OSC: AZUL_OSC,
    NARANJA:  RGBColor(0x7E, 0x44, 0x0C),
}

for col_i, details in enumerate(D_COLS):
    # flecha vertical desde caja principal hacia detalles
    arrow_v(slide,
            xs[col_i] + BW / 2 - Inches(0.1),
            BY + BH,
            Inches(0.6),
            color=list(ALPHA.values())[col_i])

    for row_i, (detail, color) in enumerate(details):
        dy = DY0 + row_i * (DH + Inches(0.12))
        bg = ALPHA[color]
        tc = TEXT_C[color]
        rect(slide, xs[col_i], dy, BW, DH, fill=bg, line=color)
        txt(slide, detail,
            xs[col_i] + Inches(0.1), dy + Inches(0.06),
            BW - Inches(0.2), DH - Inches(0.1),
            size=10, color=tc, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════════════════
# LEYENDA de salidas
# ══════════════════════════════════════════════════════════════════════════════
leg_y = Inches(6.7)
leg_items = [
    ("CSV preprocesado + alias", AZUL_MED),
    ("CSV procesado (formato largo)", AZUL_OSC),
    ("Tablas PostgreSQL", NARANJA),
]
leg_x = Inches(1.5)
txt(slide, "Salidas:", leg_x - Inches(0.8), leg_y, Inches(0.75), Inches(0.4),
    size=11, bold=True, color=AZUL_OSC)
for i, (label, color) in enumerate(leg_items):
    lx = leg_x + i * Inches(3.5)
    rect(slide, lx, leg_y + Inches(0.07), Inches(0.25), Inches(0.25), fill=color)
    txt(slide, label, lx + Inches(0.32), leg_y + Inches(0.04),
        Inches(3.1), Inches(0.35), size=11, color=NEGRO)

out = "C3_Procesador_de_datos.pptx"
prs.save(out)
print(f"Guardado: {out}")
